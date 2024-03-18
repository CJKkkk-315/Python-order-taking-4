import os
from pptx import Presentation
from PIL import Image
import io
# PPT文件路径
ppt_path = '老38-75第2列汇总1193例.pptx'

# 保存图片的文件夹路径
save_dir = ppt_path.split('.')[0]

# 创建目录
os.makedirs(save_dir, exist_ok=True)

# 加载PPT文件
presentation = Presentation(ppt_path)

# 循环遍历所有幻灯片
for slide_no, slide in enumerate(presentation.slides):

    # 循环遍历幻灯片中的所有形状
    for shape in slide.shapes:

        # 检查形状是否为图片
        if hasattr(shape, "image"):

            # 获取图片数据
            image = shape.image

            # 使用PIL将图片数据转换为图像
            pil_image = Image.open(io.BytesIO(image.blob))

            # 保存图片到指定文件夹，文件名为"slide_{slide_no}_{image_no}.png"
            image_path = os.path.join(save_dir, f'slide_{slide_no}_image.png')
            pil_image.save(image_path)
            print(f'Saved {image_path}')
